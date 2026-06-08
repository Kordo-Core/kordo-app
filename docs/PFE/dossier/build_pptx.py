# -*- coding: utf-8 -*-
"""Genere kordo-presentation.pptx (20 slides, editable PowerPoint)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---- palette ----
GREEN      = RGBColor(0x3f, 0x8f, 0x6b)
GREEN_DARK = RGBColor(0x2f, 0x6e, 0x52)
GREEN_SOFT = RGBColor(0xea, 0xf3, 0xee)
ORANGE     = RGBColor(0xe0, 0x90, 0x2f)
ORANGE_D   = RGBColor(0xb9, 0x71, 0x1a)
ORANGE_SOFT= RGBColor(0xfb, 0xf2, 0xe3)
DARK       = RGBColor(0x23, 0x23, 0x23)
MUTED      = RGBColor(0x6b, 0x6b, 0x6b)
LINE       = RGBColor(0xe2, 0xe2, 0xde)
BLUE       = RGBColor(0x3f, 0x78, 0xb5)
BLUE_SOFT  = RGBColor(0xe8, 0xf0, 0xf8)
PURPLE     = RGBColor(0x7e, 0x57, 0xa8)
PURPLE_SOFT= RGBColor(0xf1, 0xea, 0xf7)
RED        = RGBColor(0xc0, 0x56, 0x3b)
WHITE      = RGBColor(0xff, 0xff, 0xff)
GREY_OUT   = RGBColor(0x8a, 0x8a, 0x8a)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


def slide():
    return prs.slides.add_slide(BLANK)


def no_shadow(shp):
    shp.shadow.inherit = False


def bg(s, color):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SW), Inches(SH))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); no_shadow(r)
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return r


def tx(s, l, t, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT,
       anchor=MSO_ANCHOR.TOP, font=FONT, italic=False, spacing=1.0):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.name = font; f.color.rgb = color
    return tb


def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0, rounded=True, radius=0.08):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if rounded:
        try: sp.adjustments[0] = radius
        except Exception: pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    no_shadow(sp)
    return sp


def set_text(sp, text, size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE, font=FONT, spacing=1.0):
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.name = font; f.color.rgb = color


def header(s, kicker, title, title_color=DARK):
    tx(s, 0.7, 0.5, 11.9, 0.4, kicker.upper(), size=13, color=ORANGE, bold=True)
    tx(s, 0.7, 0.85, 11.9, 0.8, title, size=30, color=title_color, bold=True)
    rule = rect(s, 0.72, 1.62, 0.75, 0.06, fill=GREEN, rounded=False)
    rect(s, 1.17, 1.62, 0.3, 0.06, fill=ORANGE, rounded=False)
    return rule


def subtitle(s, text, y=1.85):
    tx(s, 0.7, y, 11.9, 0.5, text, size=14, color=MUTED)


def footer(s, n):
    tx(s, 0.7, 7.0, 6, 0.3, "Kordo · Soutenance PFE 2025–2026", size=8.5, color=MUTED)
    tx(s, 7.0, 7.0, 5.63, 0.3, f"{n} / 20", size=8.5, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)


def placeholder(s, l, t, w, h, title, desc=""):
    sp = rect(s, l, t, w, h, fill=RGBColor(0xfa, 0xfa, 0xf8), line=RGBColor(0xc4, 0xc4, 0xbe), line_w=1.5)
    set_text(sp, ("🖼  " + title + ("\n" + desc if desc else "")), size=11,
             color=RGBColor(0x9a, 0x9a, 0x93), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return sp


# =================== SLIDE 1 — COVER ===================
s = slide(); bg(s, GREEN_DARK)
rect(s, 0, 0, SW, SH, fill=GREEN, rounded=False).fill.fore_color.rgb  # base
# gradient-ish: overlay lighter rectangle on right
ov = rect(s, 7.5, 0, SW-7.5, SH, fill=RGBColor(0x4b, 0xa0, 0x7b), rounded=False)
tx(s, 0.9, 2.05, 11, 1.6, "Kordo.", size=92, color=WHITE, bold=True)
tx(s, 0.95, 3.55, 11, 0.6, "Réseau social d'escalade indoor", size=20, color=WHITE)
rect(s, 0.97, 4.25, 1.25, 0.06, fill=ORANGE, rounded=False)
tx(s, 0.95, 4.6, 11, 0.4, "Soutenance · Projet de Fin d'Études · 2025–2026", size=15, color=WHITE, bold=True)
tx(s, 0.95, 5.15, 11.5, 0.4, "Jacinto Valentino — M2 Développement, cursus Tech Lead · ECV Paris", size=14, color=WHITE)
tx(s, 6.8, 6.9, 5.8, 0.3, "🖼  Fond : photo d'ambiance bloc indoor (optionnel)", size=9, color=RGBColor(0xcf,0xe6,0xda), italic=True, align=PP_ALIGN.RIGHT)

# =================== SLIDE 2 — QUI SUIS-JE ===================
s = slide(); header(s, "Présentation", "Qui suis-je ?"); footer(s, 2)
placeholder(s, 0.7, 2.1, 5.2, 4.4, "Photo de toi à insérer", "Portrait ou photo en salle de bloc")
items = [
    "Jacinto Valentino",
    "Étudiant M2 Développement à l'ECV Paris — cursus Tech Lead",
    "Grimpeur depuis 2020 (post-Covid)",
    "Passionné d'escalade indoor et de produit mobile",
    "Kordo est né de ma propre pratique et de ses frustrations",
]
y = 2.5
for it in items:
    tx(s, 6.5, y, 0.4, 0.4, "▹", size=15, color=GREEN, bold=True)
    tx(s, 6.85, y, 5.6, 0.7, it, size=14.5, color=DARK, spacing=1.05)
    y += 0.72

# =================== SLIDE 3 — PROBLEME ===================
s = slide(); header(s, "Le constat", "Le problème"); footer(s, 3)
subtitle(s, "Un grimpeur jongle entre 4 outils — aucun ne réunit tout.")
cards = [
    ("📸  Instagram", "Pour partager ses sessions… mais aucune donnée d'escalade."),
    ("📊  Un tableur", "Pour suivre sa progression et ses blocs validés, à la main."),
    ("💬  WhatsApp", "Pour trouver un partenaire de grimpe au dernier moment."),
    ("📋  Tableau de la salle", "Pour repérer les blocs, les grades et les ouvreurs."),
]
cw = 2.78; gap = 0.18; x = 0.7
for title, desc in cards:
    c = rect(s, x, 2.5, cw, 2.3, fill=WHITE, line=LINE, line_w=1.0)
    tx(s, x+0.2, 2.7, cw-0.4, 0.5, title, size=14, color=GREEN_DARK, bold=True)
    tx(s, x+0.2, 3.3, cw-0.4, 1.3, desc, size=11, color=MUTED, spacing=1.1)
    x += cw + gap
band = rect(s, 0.7, 5.35, 11.93, 0.9, fill=ORANGE_SOFT)
set_text(band, "Résultat : une expérience fragmentée et des données perdues à chaque session.",
         size=13, color=ORANGE_D, bold=True, align=PP_ALIGN.CENTER)

# =================== SLIDE 4 — SOLUTION ===================
s = slide(); header(s, "La solution", "Kordo"); footer(s, 4)
tx(s, 0.7, 2.1, 6.0, 1.5,
   "La première plateforme indoor-first qui réunit réseau social, suivi de progression et exploration des salles dans un seul produit.",
   size=15, color=DARK, spacing=1.25)
obj = rect(s, 0.7, 4.0, 6.0, 2.2, fill=GREEN_SOFT)
tx(s, 0.95, 4.25, 5.5, 0.5, "🎯  Objectif du projet", size=14, color=GREEN_DARK, bold=True)
tx(s, 0.95, 4.85, 5.5, 1.2,
   "Valider l'usage auprès de 3 salles pilotes et 100+ grimpeurs, avec une rétention J7 > 40 % comme signal de product-market fit.",
   size=12, color=DARK, spacing=1.2)
placeholder(s, 7.2, 2.1, 5.43, 4.4, "Capture d'écran à insérer", "Écran d'accueil / feed de l'app Kordo")

# =================== SLIDE 5 — PRODUIT ===================
s = slide(); header(s, "Le produit", "Kordo en un coup d'œil"); footer(s, 5)
b1 = rect(s, 0.7, 2.1, 6.4, 1.85, fill=GREEN_SOFT)
tx(s, 0.95, 2.3, 6.0, 0.5, "4 types de publication", size=14, color=GREEN_DARK, bold=True)
tx(s, 0.95, 2.85, 6.0, 1.0, "Post (photo) · Message (texte) · Activité (séance : salle, blocs validés, vidéos) · Now (présence live)",
   size=11.5, color=DARK, spacing=1.15)
b2 = rect(s, 0.7, 4.15, 6.4, 1.95, fill=WHITE, line=LINE)
tx(s, 0.95, 4.35, 6.0, 0.5, "Indoor — le cœur de l'app", size=14, color=GREEN_DARK, bold=True)
tx(s, 0.95, 4.9, 6.0, 1.1, "Salles + filtres (favoris, populaires, récentes) · Blocs & validation · Vidéos de méthode · Classement par salle",
   size=11.5, color=DARK, spacing=1.15)
placeholder(s, 7.3, 2.1, 2.6, 4.0, "Écran Feed", "à insérer")
placeholder(s, 10.03, 2.1, 2.6, 4.0, "Fiche bloc + vidéo", "à insérer")

# =================== SLIDE 6 — MARCHE 1 ===================
s = slide(); header(s, "Étude de marché · 1/3", "Un marché en pleine croissance"); footer(s, 6)
subtitle(s, "Porté par les Jeux Olympiques de Paris 2024, l'escalade indoor explose.")
stats = [
    ("1,2 M", "pratiquants en France", GREEN_SOFT, GREEN_DARK),
    ("8 M+", "pratiquants en Europe", BLUE_SOFT, BLUE),
    ("600+", "salles indoor en France", ORANGE_SOFT, ORANGE_D),
    ("28 ans", "âge médian (Gen Z ↑)", ORANGE_SOFT, ORANGE_D),
    ("+12–15 %", "croissance par an", GREEN_SOFT, GREEN_DARK),
    ("+25 %", "d'inscriptions post-JO 2024", BLUE_SOFT, BLUE),
]
cw = 3.85; ch = 1.7; gx = 0.7; gy = 2.55
for i, (n, l, fill, col) in enumerate(stats):
    col_i = i % 3; row_i = i // 3
    x = gx + col_i * (cw + 0.19); y = gy + row_i * (ch + 0.2)
    rect(s, x, y, cw, ch, fill=fill)
    tx(s, x+0.25, y+0.25, cw-0.5, 0.7, n, size=27, color=col, bold=True)
    tx(s, x+0.25, y+1.05, cw-0.5, 0.5, l, size=11, color=MUTED)

# =================== SLIDE 7 — TAM SAM SOM ===================
s = slide(); header(s, "Étude de marché · 2/3", "Marché adressable"); footer(s, 7)
subtitle(s, "Mesuré en nombre de grimpeurs.")
# concentric circles
cx = 3.3
oval1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(2.55), Inches(3.6), Inches(3.6))
oval1.fill.solid(); oval1.fill.fore_color.rgb = GREEN_SOFT; oval1.line.fill.background(); no_shadow(oval1)
oval2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx+0.55), Inches(3.1), Inches(2.5), Inches(2.5))
oval2.fill.solid(); oval2.fill.fore_color.rgb = RGBColor(0xcf,0xe6,0xda); oval2.line.fill.background(); no_shadow(oval2)
oval3 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx+1.07), Inches(3.62), Inches(1.45), Inches(1.45))
oval3.fill.solid(); oval3.fill.fore_color.rgb = GREEN; oval3.line.fill.background(); no_shadow(oval3)
tx(s, cx, 2.65, 3.6, 0.3, "TAM", size=10, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)
tx(s, cx+0.55, 3.2, 2.5, 0.3, "SAM", size=10, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)
tx(s, cx+1.07, 4.15, 1.45, 0.3, "SOM", size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
# legend
lx = 7.6
legend = [
    ("≈ 44,5 M", "TAM", "Grimpeurs dans le monde, toutes pratiques"),
    ("≈ 8 M", "SAM", "Grimpeurs indoor en Europe équipés smartphone"),
    ("≈ 75 000", "SOM", "France + Benelux, an 1–3, pénétration ~5 %"),
]
yy = 2.75
for n, t, d in legend:
    tx(s, lx, yy, 4.8, 0.5, n + "   " + t, size=16, color=GREEN_DARK, bold=True)
    tx(s, lx, yy+0.5, 4.8, 0.5, d, size=11, color=MUTED)
    yy += 1.25

# =================== SLIDE 8 — CONCURRENCE ===================
s = slide(); header(s, "Étude de marché · 3/3", "Concurrence"); footer(s, 8)
subtitle(s, "Chacun couvre une brique — aucun ne réunit le tout pour l'indoor.")
comp = [
    ("Arkose+", "La salle", "App d'enseigne : résa, infos. Mono-enseigne, pas de social ouvert."),
    ("Strava", "Le social sportif", "Feed & gamification puissants. Mais pas d'escalade, ni blocs ni salles."),
    ("Instagram", "Le partage", "Audience massive. Aucune donnée grimpe structurée ni classement."),
    ("8a.nu", "Le carnet de blocs", "Logbook & base de voies. Orienté outdoor, UX datée, social faible."),
]
cw = 2.85; x = 0.7
for nm, role, ms in comp:
    head = rect(s, x, 2.6, cw, 0.6, fill=DARK)
    set_text(head, nm, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    body = rect(s, x, 3.2, cw, 1.75, fill=WHITE, line=LINE)
    tx(s, x+0.18, 3.35, cw-0.36, 0.4, role, size=11, color=ORANGE_D, bold=True)
    tx(s, x+0.18, 3.8, cw-0.36, 1.1, ms, size=10, color=MUTED, spacing=1.1)
    x += cw + 0.19
band = rect(s, 0.7, 5.2, 11.93, 0.85, fill=GREEN_SOFT)
set_text(band, "Kordo = la seule plateforme à réunir salle indoor + progression + social en un produit.",
         size=13, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

# =================== SLIDE 9 — BUSINESS MODEL ===================
s = slide(); header(s, "Business Plan · 1/2", "Modèle économique"); footer(s, 9)
subtitle(s, "Un modèle hybride B2B2C aligné sur l'usage réel.")
rev = [
    ("B2B · Salles", "Setup Fee", "500 – 2 000 €", "Frais unique d'intégration, configuration des blocs & formation."),
    ("B2B · Salles", "SaaS récurrent", "1 – 3 €/user/mois", "Abonnement selon le nombre d'utilisateurs actifs de la salle."),
    ("B2C · Grimpeurs", "Premium", "4,99 €/mois", "Stats avancées, badges exclusifs, personnalisation."),
    ("B2C · Grimpeurs", "Freemium", "Gratuit", "Feed, profil, blocs — acquisition & masse critique."),
]
cw = 2.85; x = 0.7
for tag, h, pr, p in rev:
    c = rect(s, x, 2.5, cw, 2.6, fill=WHITE, line=LINE)
    tx(s, x+0.18, 2.65, cw-0.36, 0.35, tag.upper(), size=9.5, color=ORANGE, bold=True)
    tx(s, x+0.18, 3.05, cw-0.36, 0.45, h, size=14, color=DARK, bold=True)
    tx(s, x+0.18, 3.55, cw-0.36, 0.4, pr, size=12, color=GREEN_DARK, bold=True)
    tx(s, x+0.18, 4.05, cw-0.36, 1.0, p, size=9.5, color=MUTED, spacing=1.1)
    x += cw + 0.19
band = rect(s, 0.7, 5.45, 11.93, 0.8, fill=GREEN_SOFT)
set_text(band, "Revenus salles prévisibles + conversion freemium → premium côté grimpeurs.",
         size=13, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

# =================== SLIDE 10 — GTM + PROJECTIONS ===================
s = slide(); header(s, "Business Plan · 2/2", "Go-to-market & projections"); footer(s, 10)
phases = [
    ("Pilote", "T3 2026", "3 salles partenaires\nMVP déployé\n500 utilisateurs", GREEN),
    ("Lancement", "T4 2026", "App Store + Play Store\nInfluenceurs & PR\n3 000 téléchargements", ORANGE),
    ("Croissance", "2027", "20 salles\n5 000 actifs/mois\nPremium lancé", GREEN),
    ("Expansion", "2028", "BE, CH, ES, DE\n50+ salles · 20 000 actifs\nLevée Série A", GREEN),
]
cw = 2.85; x = 0.7
for ph, pw, body, col in phases:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+cw/2-0.08), Inches(2.0), Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background(); no_shadow(dot)
    box = rect(s, x, 2.35, cw, 1.95, fill=WHITE, line=LINE)
    tx(s, x+0.18, 2.5, cw-0.36, 0.4, ph, size=13, color=DARK, bold=True)
    tx(s, x+0.18, 2.9, cw-0.36, 0.3, pw, size=9.5, color=ORANGE, bold=True)
    tx(s, x+0.18, 3.25, cw-0.36, 1.0, body, size=9.5, color=MUTED, spacing=1.15)
    x += cw + 0.19
proj = [
    ("~ 9 k€", "Revenu année 1 — amorçage", GREEN_SOFT, GREEN_DARK),
    ("~ 84 k€", "Revenu année 2 — seuil de rentabilité", ORANGE_SOFT, ORANGE_D),
    ("~ 300 k€", "Revenu année 3 — croissance rentable", BLUE_SOFT, BLUE),
]
cw2 = 3.85; x = 0.7
for n, l, fill, col in proj:
    rect(s, x, 4.65, cw2, 1.55, fill=fill)
    tx(s, x+0.25, 4.9, cw2-0.5, 0.6, n, size=24, color=col, bold=True)
    tx(s, x+0.25, 5.6, cw2-0.5, 0.5, l, size=10.5, color=MUTED)
    x += cw2 + 0.19

# =================== SLIDE 11 — SWOT ===================
s = slide(); header(s, "Stratégie", "Analyse SWOT"); footer(s, 11)
swot = [
    ("Forces", GREEN_SOFT, GREEN_DARK, ["1er réseau social escalade indoor", "Approche B2B2C (salles + grimpeurs)", "Vidéos de méthode sur les blocs — unique", "Stack moderne & design system Kordo-UI"]),
    ("Faiblesses", ORANGE_SOFT, ORANGE_D, ["Équipe réduite, 0 traction commerciale", "Pas de financement initial", "Dépendance à l'adoption des salles", "Notoriété nulle au démarrage"]),
    ("Opportunités", BLUE_SOFT, BLUE, ["Marché +12 %/an, boost post-JO durable", "Gen Z très tech-native", "Peu de concurrents sur l'indoor social", "SaaS B2B scalable (600+ salles FR)"]),
    ("Menaces", PURPLE_SOFT, PURPLE, ["Strava / Arkose comblant la brique manquante", "Instagram & TikTok captent l'attention", "Cycles d'acquisition B2B longs", "RGPD sur les données sportives"]),
]
qw = 5.87; qh = 2.18
positions = [(0.7, 2.1), (6.76, 2.1), (0.7, 4.42), (6.76, 4.42)]
for (qx, qy), (title, fill, col, items) in zip(positions, swot):
    rect(s, qx, qy, qw, qh, fill=fill)
    tx(s, qx+0.25, qy+0.15, qw-0.5, 0.4, title, size=13, color=col, bold=True)
    tx(s, qx+0.25, qy+0.6, qw-0.5, qh-0.7, "\n".join("•  " + it for it in items), size=10.5, color=DARK, spacing=1.25)

# =================== SLIDE 12 — ROADMAP ===================
s = slide(); header(s, "Pilotage projet", "Roadmap — vue d'ensemble"); footer(s, 12)
road = [
    ("Sprint 1–2", "Mars 2026 · Terminés", "Init projet\nDesign System\nKordo-UI stable", GREEN),
    ("Sprint 3–4", "Mar–Avr 2026", "Auth · Profil\nFeed & publications\nNotifications", ORANGE),
    ("Sprint 5–6", "Avr–Mai 2026", "Indoor · blocs · vidéos\nClassements\nMessagerie · paramètres", ORANGE),
    ("MVP Launch", "Juin 2026", "App Store + Play Store\n3 salles pilotes\nMonitoring & RGPD", RED),
]
cw = 2.85; x = 0.7
for ph, pw, body, col in road:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+cw/2-0.09), Inches(2.4), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background(); no_shadow(dot)
    box = rect(s, x, 2.8, cw, 2.5, fill=WHITE, line=LINE)
    tx(s, x+0.18, 3.0, cw-0.36, 0.4, ph, size=14, color=DARK, bold=True)
    tx(s, x+0.18, 3.45, cw-0.36, 0.3, pw, size=9.5, color=ORANGE, bold=True)
    tx(s, x+0.18, 3.85, cw-0.36, 1.3, body, size=10.5, color=MUTED, spacing=1.2)
    x += cw + 0.19

# =================== SLIDE 13 — AGILE ===================
s = slide(); header(s, "Pilotage projet", "Méthode agile"); footer(s, 13)
flow = [("Backlog", "Épics & user stories", GREEN), ("Sprint", "Itérations de 2 semaines", ORANGE), ("Tâches", "Découpage & réalisation", BLUE)]
bw = 3.4; x = 0.9
for i, (h, p, col) in enumerate(flow):
    b = rect(s, x, 2.3, bw, 1.4, fill=col)
    set_text(b, h + "\n" + p, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < 2:
        tx(s, x+bw+0.05, 2.55, 0.6, 0.9, "→", size=26, color=MUTED, align=PP_ALIGN.CENTER)
    x += bw + 0.65
tx(s, 0.7, 4.1, 11.9, 0.5, "Priorisation et estimation structurées :", size=14, color=MUTED)
chips = ["MoSCoW (Critique → Faible)", "Estimation Fibonacci (1·2·3·5·8·13)", "Vélocité ~30 pts / sprint", "25 user stories · 124 points", "Dépendances explicites"]
x = 0.7; y = 4.7
for c in chips:
    w = 0.28 + len(c) * 0.092
    ch = rect(s, x, y, w, 0.55, fill=GREEN_SOFT, radius=0.5)
    set_text(ch, c, size=11.5, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)
    x += w + 0.2
    if x > 10.5:
        x = 0.7; y += 0.75

# =================== SLIDE 14 — STORY MAPPING ===================
s = slide(); header(s, "Pilotage projet", "Story Mapping"); footer(s, 14)
subtitle(s, "Parcours utilisateur (colonnes) × priorité (lignes). La ligne orange = frontière MVP.")
cols = [("ONBOARDING", GREEN), ("INDOOR", ORANGE), ("FEED & SOCIAL", RGBColor(0xc9,0x77,0x2f)), ("PROFIL", PURPLE), ("MESSAGES", BLUE)]
grid = [
    ["Inscription / Connexion", "Liste salles + filtres", "Fil d'actualité", "Profil public", "Notifications"],
    ["OAuth Google / Apple", "Blocs + validation", "Post·Message·Activité·Now", "Modifier profil + avatar", "Messagerie 1:1"],
    ["", "Vidéos méthode · Classement", "Like & commentaire", "Follow / Abonnements", ""],
]
cw = 2.3; gap = 0.12; x0 = 0.72
for i, (name, col) in enumerate(cols):
    x = x0 + i * (cw + gap)
    h = rect(s, x, 2.45, cw, 0.45, fill=col)
    set_text(h, name, size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
for r, row in enumerate(grid):
    for i, cell in enumerate(row):
        if not cell: continue
        x = x0 + i * (cw + gap); y = 3.0 + r * 0.52
        c = rect(s, x, y, cw, 0.45, fill=GREEN_SOFT)
        set_text(c, cell, size=8.5, color=RGBColor(0x22,0x33,0x44), align=PP_ALIGN.CENTER)
rect(s, 0.72, 4.62, 11.86, 0.05, fill=ORANGE, rounded=False)
postmvp = ["Suppression compte", "Check-in · Détection géoloc", "Now auto · Résumé hebdo", "Graphiques · Premium", "Appels · Événements"]
for i, cell in enumerate(postmvp):
    x = x0 + i * (cw + gap)
    c = rect(s, x, 4.78, cw, 0.45, fill=RGBColor(0xef,0xef,0xef))
    set_text(c, cell, size=8.5, color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.CENTER)
tx(s, 0.7, 5.4, 11.9, 0.3, "Vert = MVP (S3–S6) · Gris = post-MVP", size=9, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# =================== SLIDE 15 — MVP ===================
s = slide(); header(s, "Le périmètre", "Le MVP : choisi vs reporté"); footer(s, 15)
inn = ["Auth email + Google / Apple", "Profil public + modification & avatar", "Feed + 4 publications (Now manuel)", "Like & commentaire", "Follow / unfollow", "Notifications + messagerie (sans appel)", "Salles + filtres · blocs + validation", "Vidéos de méthode · classement salle"]
out = ["Check-in dématérialisé + Now auto", "Détection grimpeurs (Meet, géoloc)", "Graphiques de progression", "Badges & trophées", "Compte Premium B2C", "Création d'événements · résumé hebdo", "Mode sombre · multilingue · appels", "Dashboard gérant (B2B)"]
# IN
hin = rect(s, 0.7, 2.1, 5.87, 0.55, fill=GREEN)
set_text(hin, "✓  Dans le MVP", size=14, color=WHITE, bold=True)
bin_ = rect(s, 0.7, 2.65, 5.87, 3.0, fill=WHITE, line=LINE)
tx(s, 0.95, 2.8, 5.4, 2.8, "\n".join("✓  " + i for i in inn), size=11.5, color=DARK, spacing=1.35)
# OUT
hout = rect(s, 6.76, 2.1, 5.87, 0.55, fill=GREY_OUT)
set_text(hout, "✗  Reporté (post-MVP)", size=14, color=WHITE, bold=True)
bout = rect(s, 6.76, 2.65, 5.87, 3.0, fill=WHITE, line=LINE)
tx(s, 7.0, 2.8, 5.4, 2.8, "\n".join("✗  " + i for i in out), size=11.5, color=MUTED, spacing=1.35)
band = rect(s, 0.7, 5.85, 11.93, 0.8, fill=ORANGE_SOFT)
set_text(band, "Le vrai défi : couper des features d'un projet auquel on croit, pour livrer l'essentiel.",
         size=12.5, color=ORANGE_D, bold=True, align=PP_ALIGN.CENTER)

# =================== SLIDE 16 — STACK ===================
s = slide(); header(s, "Développement · 1/2", "Stack technique"); footer(s, 16)
techs = [
    ("React Native + Expo", "Application mobile cross-platform iOS & Android", GREEN),
    ("TypeScript", "Typage strict de bout en bout", BLUE),
    ("Supabase", "Auth · PostgreSQL · Realtime · Storage", ORANGE),
    ("expo-router", "Navigation par fichiers", PURPLE),
    ("Emotion", "Styling cross-platform (native + web)", ORANGE),
    ("pnpm monorepo", "core · kordo-ui · frontend · backend", BLUE),
    ("Kordo-UI", "Design system atoms → organisms, Storybook", GREEN),
    ("Sentry · PostHog", "Monitoring & analytics produit", PURPLE),
]
cw = 3.85; ch = 1.05; x0 = 0.7; y0 = 2.3
for i, (h, p, col) in enumerate(techs):
    ci = i % 3; ri = i // 3
    x = x0 + ci * (cw + 0.19); y = y0 + ri * (ch + 0.22)
    rect(s, x, y, cw, ch, fill=WHITE, line=LINE)
    rect(s, x, y, 0.08, ch, fill=col, rounded=False)
    tx(s, x+0.25, y+0.13, cw-0.4, 0.4, h, size=13, color=DARK, bold=True)
    tx(s, x+0.25, y+0.55, cw-0.4, 0.45, p, size=10, color=MUTED)

# =================== SLIDE 17 — ARCHITECTURE ===================
s = slide(); header(s, "Développement · 2/2", "Architecture"); footer(s, 17)
def archbox(s, l, t, w, h, color, title, sub):
    b = rect(s, l, t, w, h, fill=color)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sub; r2.font.size = Pt(9.5); r2.font.color.rgb = WHITE; r2.font.name = FONT
archbox(s, 4.67, 2.0, 4.0, 0.95, GREEN, "📱  Application mobile", "React Native + Expo · expo-router")
tx(s, 4.67, 3.05, 4.0, 0.4, "↓  API REST · Realtime", size=14, color=MUTED, align=PP_ALIGN.CENTER)
archbox(s, 4.67, 3.5, 4.0, 0.95, BLUE, "⚙  Supabase", "Auth · Realtime · Storage · Edge Functions")
tx(s, 4.67, 4.55, 4.0, 0.4, "↓", size=16, color=MUTED, align=PP_ALIGN.CENTER)
archbox(s, 4.67, 4.95, 4.0, 0.95, DARK, "🗄  PostgreSQL", "Users · Gyms · Boulders · Posts · Messages")
mono = [("core", "types domaine"), ("kordo-ui", "design system"), ("frontend", "app Expo"), ("backend", "API")]
mw = 2.85; x = 0.7
for nm, d in mono:
    m = rect(s, x, 6.15, mw, 0.75, fill=GREEN_SOFT)
    set_text(m, nm + "\n" + d, size=10.5, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)
    x += mw + 0.19

# =================== SLIDE 18 — DIFFICULTES ===================
s = slide(); header(s, "Retour d'expérience", "Difficultés rencontrées"); footer(s, 18)
diffs = [
    ("1", "Les contraintes de temps", "Mener de front conception, design, architecture et développement en solo impose des arbitrages permanents sur le périmètre."),
    ("2", "Réduire le périmètre du MVP", "Vouloir tout mettre dans le MVP. Couper des fonctionnalités d'un projet auquel on croit est difficile — apprendre à prioriser (MoSCoW, frontière MVP) a été clé."),
    ("3", "Un business model réaliste & rentable", "Construire un modèle viable : équilibrer la valeur B2B (salles) et B2C (grimpeurs) pour quelque chose de crédible et rentable."),
]
y = 2.2
for num, h, p in diffs:
    box = rect(s, 0.7, y, 11.93, 1.35, fill=WHITE, line=LINE)
    tx(s, 0.95, y+0.3, 0.7, 0.8, num, size=26, color=ORANGE, bold=True)
    tx(s, 1.75, y+0.2, 10.6, 0.45, h, size=14, color=DARK, bold=True)
    tx(s, 1.75, y+0.66, 10.6, 0.6, p, size=11, color=MUTED, spacing=1.1)
    y += 1.52

# =================== SLIDE 19 — DEMO ===================
s = slide(); header(s, "Place à la pratique", "Démo du MVP"); footer(s, 19)
placeholder(s, 0.7, 2.2, 2.85, 3.8, "Écran 1", "Feed / publication")
placeholder(s, 3.7, 2.2, 2.85, 3.8, "Écran 2", "Indoor / fiche bloc")
card = rect(s, 6.85, 2.2, 5.78, 1.9, fill=GREEN_SOFT)
tx(s, 7.1, 2.4, 5.3, 0.5, "🎬  Démonstration live", size=14, color=GREEN_DARK, bold=True)
tx(s, 7.1, 2.95, 5.3, 1.1, "Parcours : inscription → exploration d'une salle → validation d'un bloc → publication d'une activité → feed social.",
   size=12, color=DARK, spacing=1.25)
placeholder(s, 6.85, 4.3, 5.78, 1.7, "QR code / lien démo (optionnel)")

# =================== SLIDE 20 — MERCI ===================
s = slide(); bg(s, GREEN_DARK)
rect(s, 0, 0, SW, SH, fill=GREEN, rounded=False)
rect(s, 7.5, 0, SW-7.5, SH, fill=RGBColor(0x37,0x83,0x62), rounded=False)
tx(s, 0, 2.3, SW, 1.3, "Merci !", size=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tx(s, 0, 3.55, SW, 0.7, "Questions ?", size=22, color=WHITE, align=PP_ALIGN.CENTER)
tx(s, 0, 4.7, SW, 0.5, "Jacinto Valentino · M2 Développement — ECV Paris · 2025–2026", size=14, color=RGBColor(0xe6,0xf2,0xec), align=PP_ALIGN.CENTER)
tx(s, 0, 5.2, SW, 0.8, "Kordo.", size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "kordo-presentation.pptx")
prs.save(out)
print("PPTX OK:", out, "·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
